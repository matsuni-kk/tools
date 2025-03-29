#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Context Extracter UI (Modern Design)

指定したフォルダ名を含むすべてのフォルダとそのサブフォルダ内のファイル内容を抽出。
すべてのファイル内容を1つのMarkdownファイルに集約する。
モダンなUIインターフェースを備え、結果をクリップボードにコピーできる。

使用例:
        python context_extracter_ui.py
"""

import os
import sys
import platform
import datetime
import time
import argparse
from pathlib import Path
import logging
import re
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import threading
import queue
import csv # csv モジュールをインポート

# pywin32 のインポート (Windowsのみ)
import platform # platform をインポート (重複インポート防止のため確認)
HAS_PYWIN32 = False
if platform.system() == "Windows":
    try:
        import win32com.client
        import pythoncom # COM初期化に必要
        HAS_PYWIN32 = True
    except ImportError:
        print("警告: 'pywin32' ライブラリが見つかりません。.ppt ファイルの処理はスキップされます。Windows環境の場合 `pip install pywin32` でインストールしてください。", file=sys.stderr)
else:
    # Windows以外の場合も警告を表示
    print("情報: Windows以外の環境では .ppt ファイルの処理はサポートされません。", file=sys.stderr)

# 必要なライブラリをインポート
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    print("警告: 'openpyxl' ライブラリが見つかりません。Excel (.xlsx) ファイルの処理はスキップされます。`pip install openpyxl` でインストールしてください。", file=sys.stderr)

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("警告: 'python-pptx' ライブラリが見つかりません。PowerPoint (.pptx) ファイルの処理はスキップされます。`pip install python-pptx` でインストールしてください。", file=sys.stderr)

try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False
    print("警告: 'PyPDF2' ライブラリが見つかりません。PDF (.pdf) ファイルの処理はスキップされます。`pip install PyPDF2` でインストールしてください。", file=sys.stderr)

# output ディレクトリの存在確認または作成
# __file__ から絶対パスを取得し、その親ディレクトリを基準にする
try:
    # sys.argv[0]を使用してスクリプトの絶対パスを取得（MacとWindowsの両方で動作）
    script_dir = Path(os.path.abspath(sys.argv[0])).parent
    output_dir = script_dir / "output"
    output_dir.mkdir(exist_ok=True)
    log_file_path = output_dir / "context_extracter_ui.log"
except Exception as e:
    # outputディレクトリ作成やログファイルパス設定に失敗した場合
    # スクリプトと同じディレクトリにログを出力するなどのフォールバック処理
    script_dir = Path(__file__).resolve().parent
    log_file_path = script_dir / "context_extracter_ui.fallback.log"
    print(f"警告: outputディレクトリの準備に失敗しました。ログは {log_file_path} に出力されます。エラー: {e}", file=sys.stderr)
    # output_dir が未定義の場合に備え、スクリプトディレクトリを割り当てておく
    if 'output_dir' not in locals():
        output_dir = script_dir

# ロガーの設定
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(log_file_path, mode="w", encoding="utf-8"), # 上書きモードに設定
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

logger.info(f"スクリプトディレクトリ: {script_dir}")
logger.info(f"出力ディレクトリ: {output_dir}")
logger.info(f"ログファイルパス: {log_file_path}")

def get_bms_root():
    """
    BMSのルートディレクトリを取得する。環境に合わせて調整し、存在確認を行う。
    """
    system = platform.system()
    bms_root = None
    if system == "Windows":
        # Windowsの場合は D:/BMS を想定
        bms_root = Path("D:/BMS")
    elif system == "Darwin": # macOS の正式名は Darwin
        # macOSの場合はホームディレクトリ下の Desktop/BMS を想定
        home_dir = Path.home() # ホームディレクトリを取得
        bms_root = home_dir / "Desktop" / "BMS"
        logger.info(f"macOS用BMSルートディレクトリを設定: {bms_root}")
    else: # その他のLinuxなど
        # 必要であれば他のOS用のパスも定義
        logger.warning(f"未対応のOSです: {system}. デフォルトパスは設定されません。")
        return "" # 空文字を返すか、適切なデフォルトパスを設定

    # パスが存在するか確認
    if bms_root and not bms_root.exists():
        logger.warning(f"デフォルトのBMSルートディレクトリが見つかりません: {bms_root}")
        # 存在しない場合でもパス自体は返す（ユーザーが手動で変更できるように）

    # パスが存在しても文字列として返す
    return str(bms_root) if bms_root else ""

def is_binary_file(file_path):
    """
    ファイルがバイナリかテキストかを判定
    ただし、特定の拡張子（xlsx, pptx, pdf, ppt, csv, xlsm）は内容抽出を試みるため除外
    """
    file_path_obj = Path(file_path)
    ext = file_path_obj.suffix.lower()

    # 特定の拡張子はバイナリ判定を行わず、専用の読み込み処理に任せる
    if ext in ['.xlsx', '.pptx', '.pdf', '.ppt', '.csv', '.xlsm']:
        return False # バイナリではない扱いにする（実際はバイナリだが、専用処理があるため）

    textchars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7f})
    try:
        with open(file_path, 'rb') as f:
            is_binary = bool(f.read(1024).translate(None, textchars))
        return is_binary
    except Exception as e:
        logger.warning(f"ファイル {file_path} の判定中にエラー発生: {e}")
        return True  # エラーが発生した場合もバイナリとして扱う

def is_file_in_date_range(file_path, days=None):
    """
    ファイルが指定された日数以内かどうか確認
    daysがNoneの場合は、すべてのファイルが対象
    """
    if days is None:
        return True
    
    try:
        # ファイルの最終更新日時を取得
        file_mtime = os.path.getmtime(file_path)
        now = time.time()
        
        # 日数を秒に変換
        days_in_seconds = days * 24 * 60 * 60
        
        # 指定期間内かチェック
        return (now - file_mtime) <= days_in_seconds
    except Exception as e:
        logger.warning(f"ファイル日付の確認中にエラー発生: {file_path} - {e}")
        return True  # エラーが発生した場合は含める

def match_exclude_pattern(path, exclude_patterns):
    """
    指定された除外パターンにパスが一致するかどうか確認
    いずれかのパターンに一致すればTrueを返す
    """
    if not exclude_patterns:
        return False
    
    path_lower = str(path).lower()
    for pattern in exclude_patterns:
        if pattern.lower() in path_lower:
            logger.debug(f"除外パターン '{pattern}' がパス '{path}' に一致")
            return True
    return False

def read_excel_content(file_path):
    """Excel (.xlsx) ファイルの内容を抽出"""
    if not HAS_OPENPYXL:
        return f"[ライブラリ不足のためスキップ (openpyxl): {file_path}]"
    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        content = []
        for sheet_name in workbook.sheetnames:
            content.append(f"--- シート: {sheet_name} ---")
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                row_content = []
                for cell in row:
                    if cell.value is not None:
                        row_content.append(str(cell.value))
                if row_content:
                    content.append("\t".join(row_content))
        return "\n".join(content)
    except Exception as e:
        logger.warning(f"Excelファイル {file_path} の読込中にエラー発生: {e}")
        return f"[Excelファイル読込エラー: {file_path}]"

def read_pptx_content(file_path):
    """PowerPoint (.pptx) ファイルの内容を抽出"""
    if not HAS_PPTX:
        return f"[ライブラリ不足のためスキップ (python-pptx): {file_path}]"
    try:
        prs = Presentation(file_path)
        content = []
        for i, slide in enumerate(prs.slides):
            content.append(f"--- スライド {i+1} ---")
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_texts.append(shape.text)
            # ノートも取得する場合
            if slide.has_notes_slide:
                 notes_slide = slide.notes_slide
                 notes_text = notes_slide.notes_text_frame.text
                 if notes_text:
                     content.append(f"[ノート]\n{notes_text}")

            if slide_texts:
                content.append("\n".join(slide_texts))

        return "\n".join(content)
    except Exception as e:
        logger.warning(f"PowerPointファイル {file_path} の読込中にエラー発生: {e}")
        return f"[PowerPointファイル読込エラー: {file_path}]"

def read_pdf_content(file_path):
    """PDF (.pdf) ファイルの内容を抽出"""
    if not HAS_PYPDF2:
        return f"[ライブラリ不足のためスキップ (PyPDF2): {file_path}]"
    try:
        content = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            num_pages = len(reader.pages)
            for i in range(num_pages):
                page = reader.pages[i]
                page_text = page.extract_text()
                if page_text: # テキストが存在する場合のみ追加
                     content.append(f"--- ページ {i+1} ---")
                     content.append(page_text)
        # PyPDF2がテキストを抽出できなかった場合（画像PDFなど）
        if not content:
             return f"[PDFからテキスト抽出不可（画像PDF等の可能性）: {file_path}]"

        return "\n".join(content)
    except Exception as e:
        logger.warning(f"PDFファイル {file_path} の読込中にエラー発生: {e}")
        # PyPDF2は暗号化されたPDFなどでエラーを出すことがある
        if "encrypted" in str(e).lower():
            return f"[暗号化されたPDFのため読込不可: {file_path}]"
        return f"[PDFファイル読込エラー: {file_path}]"

def read_ppt_content(file_path):
    """PowerPoint (.ppt) ファイルの内容を抽出 (Windowsのみ, 要PowerPointインストール)"""
    if platform.system() != "Windows":
        return f"[Windows環境ではないためスキップ (.ppt): {file_path}]"
    if not HAS_PYWIN32:
        return f"[ライブラリ不足のためスキップ (pywin32): {file_path}]"

    powerpoint = None
    presentation = None
    abs_path = str(Path(file_path).resolve()) # COMに渡すため絶対パスに変換

    try:
        # COMライブラリを初期化 (マルチスレッド環境で必要)
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # PowerPointを非表示で実行
        powerpoint.Visible = False
        # ファイルを開く (読み取り専用、タイトル表示なし、ウィンドウ非表示)
        presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True, Untitled=False, WithWindow=False)

        content = []
        for i, slide in enumerate(presentation.Slides):
            content.append(f"--- スライド {i+1} ---")
            slide_texts = []
            try:
                for shape in slide.Shapes:
                    # テキストフレームを持つシェイプからテキストを抽出
                    if shape.HasTextFrame:
                        if shape.TextFrame.HasText:
                            slide_texts.append(shape.TextFrame.TextRange.Text)
            except Exception as shape_e:
                 logger.warning(f".pptファイル {file_path} スライド {i+1} のシェイプ処理中にエラー: {shape_e}")
                 slide_texts.append(f"[シェイプ読込エラー: {shape_e}]")

            # ノートページのテキストを取得
            try:
                if slide.HasNotesPage:
                    # ノートは通常 Placeholders(2) にある
                    if len(slide.NotesPage.Shapes.Placeholders) >= 2:
                        notes_shape = slide.NotesPage.Shapes.Placeholders[2]
                        if notes_shape.HasTextFrame:
                            if notes_shape.TextFrame.HasText:
                                notes_text = notes_shape.TextFrame.TextRange.Text
                                if notes_text.strip(): # 空でないノートのみ追加
                                    content.append(f"[ノート]\n{notes_text}")
                    else:
                         # 稀にプレースホルダー構成が異なる場合、他のテキストを持つシェイプを探す
                         for notes_shape in slide.NotesPage.Shapes:
                              if notes_shape.HasTextFrame and notes_shape.TextFrame.HasText:
                                   notes_text = notes_shape.TextFrame.TextRange.Text
                                   if notes_text.strip():
                                        content.append(f"[ノート(代替)]\n{notes_text}")
                                        break # 最初に見つかったものを採用

            except Exception as note_e:
                 logger.warning(f".pptファイル {file_path} スライド {i+1} のノート処理中にエラー: {note_e}")
                 content.append(f"[ノート読込エラー: {note_e}]")


            if slide_texts:
                content.append("\n".join(slide_texts))

        return "\n".join(content)
    except pythoncom.com_error as e:
         # COMエラーの詳細を出力
         logger.error(f".pptファイル {file_path} の読込中にCOMエラー発生: {e}")
         hresult = getattr(e, 'hresult', 'N/A')
         desc = getattr(e, 'excepinfo', ('','','',''))[2] if getattr(e, 'excepinfo', None) else "詳細不明"
         # PowerPointがインストールされていない場合などの特定のエラーを判別
         if hresult == -2147221005: # REGDB_E_CLASSNOTREG (クラスが登録されていない)
             return f"[PowerPointがインストールされていないか、COM登録に問題があります: {file_path}]"
         return f"[.pptファイル読込COMエラー (HRESULT: {hresult}, Desc: {desc}): {file_path}]"
    except Exception as e:
        logger.warning(f".pptファイル {file_path} の読込中に予期せぬエラー発生: {e}")
        return f"[.pptファイル読込エラー: {file_path}]"
    finally:
        try:
            if presentation:
                presentation.Close()
            if powerpoint:
                powerpoint.Quit()
        except Exception as close_e:
             logger.warning(f"PowerPoint終了処理中にエラー: {close_e}")
        # COMライブラリを解放
        pythoncom.CoUninitialize()

def read_csv_content(file_path):
    """CSV (.csv) ファイルの内容を抽出 (タブ区切りも考慮し、エンコーディングを試行)"""
    # 日本語でよく使われるエンコーディングを優先的に試す
    encodings_to_try = ['utf-8', 'shift-jis', 'cp932', 'euc-jp', 'iso2022-jp', 'utf-16']
    content = f"[CSVファイル: {file_path}]\n"
    detected_encoding = None
    detected_delimiter = None

    for encoding in encodings_to_try:
        try:
            with open(file_path, 'r', encoding=encoding, newline='') as f:
                # 最初の数行を読み込んで区切り文字を推定
                sample = "".join(f.readline() for _ in range(10)) # 判定のため少し多めに読む
                if not sample:
                    return f"[空のCSVファイル: {file_path}]"

                sniffer = csv.Sniffer()
                try:
                    dialect = sniffer.sniff(sample, delimiters=',\t;|') # 一般的な区切り文字候補
                    detected_delimiter = dialect.delimiter
                except csv.Error:
                    # 自動判定失敗時はカンマをデフォルトとするが、タブの可能性も考慮
                    if '\t' in sample and ',' not in sample:
                         detected_delimiter = '\t'
                    else:
                         detected_delimiter = ',' # デフォルトはカンマ

                f.seek(0) # ファイルポインタを先頭に戻す
                # 区切り文字を確定してリーダーを作成
                reader = csv.reader(f, delimiter=detected_delimiter)
                csv_data = []
                for row in reader:
                    # 各セルを文字列に変換し、空文字はそのまま保持
                    csv_data.append(detected_delimiter.join(map(str, row)))
                detected_encoding = encoding
                # 抽出内容の前にエンコーディングと区切り文字情報を追加
                content += f"--- (エンコーディング: {detected_encoding}, 区切り文字: '{detected_delimiter}') ---\n"
                content += "\n".join(csv_data)
                return content # 読み込めたらループを抜ける
        except UnicodeDecodeError:
            continue # 次のエンコーディングを試す
        except Exception as e:
            logger.warning(f"CSVファイル {file_path} (エンコーディング: {encoding}) の読込中にエラー発生: {e}")
            # 一度エラーが出ても他のエンコーディングで成功する可能性があるのでループは継続

    # すべてのエンコーディングで失敗した場合
    logger.warning(f"CSVファイル {file_path} の適切なエンコーディングが見つかりませんでした。試行したエンコーディング: {encodings_to_try}")
    return f"[CSVファイル読込エラー(エンコーディング不明またはファイル破損): {file_path}]"

def read_file_content(file_path):
    """
    ファイルの内容を読む。エラーで対応
    バイナリファイルの場合は、[バイナリファイル]と表示
    Excel, PowerPoint, PDF, CSV は専用関数で内容を抽出
    """
    file_path_obj = Path(file_path)
    ext = file_path_obj.suffix.lower()

    # 拡張子に応じて処理を分岐
    if ext == '.xlsx' or ext == '.xlsm':
        return read_excel_content(file_path)
    elif ext == '.pptx':
        return read_pptx_content(file_path)
    elif ext == '.ppt':
        return read_ppt_content(file_path)
    elif ext == '.pdf':
        return read_pdf_content(file_path)
    elif ext == '.csv':
        return read_csv_content(file_path)
    # 既存のバイナリ判定とテキスト読み込み
    elif is_binary_file(file_path): # is_binary_fileがTrueを返すのは上記拡張子以外
        return f"[バイナリファイル: {file_path}]"
    else: # テキストファイルの場合
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # UTF-8 で失敗した場合、Shift-JIS を試す
            try:
                with open(file_path, 'r', encoding='shift-jis') as f:
                    return f.read()
            except UnicodeDecodeError:
                # Shift-JIS でも失敗した場合、CP932 を試す
                try:
                    with open(file_path, 'r', encoding='cp932') as f:
                        return f.read()
                except Exception as e:
                    logger.warning(f"ファイル {file_path} の読込中にエラー発生 (UTF-8, Shift-JIS, CP932試行後): {e}")
                    return f"[ファイル読込エラー(エンコーディング不明): {file_path}]"
            except Exception as e:
                logger.warning(f"ファイル {file_path} の読込中にエラー発生 (Shift-JIS試行後): {e}")
                return f"[ファイル読込エラー(Shift-JIS試行後): {file_path}]"
        except Exception as e:
            logger.warning(f"ファイル {file_path} の読込中にエラー発生 (UTF-8試行時): {e}")
            return f"[ファイル読込エラー: {file_path}]"

def find_folders_and_files(target_folder_path, days=None, exclude_patterns=None, status_queue=None):
    """
    指定されたフォルダパスとそのサブフォルダ内のファイルを検索
    days: 指定された日数以内のファイルのみを含める（Noneの場合はすべてのファイル）
    exclude_patterns: 除外するパターンのリスト
    status_queue: 状態更新用のキュー
    """
    found_files = [] # フォルダ情報ではなく、ファイルパスのリストに変更
    
    # target_folder_pathがPathオブジェクトの場合は文字列に変換
    if isinstance(target_folder_path, Path):
        target_folder_path = str(target_folder_path)

    log_message = f"検索開始: フォルダパス '{target_folder_path}' 内を検索"
    logger.info(log_message)
    if status_queue:
        status_queue.put(log_message)

    if days is not None:
        log_message = f"期間制限: 過去{days}日以内に更新されたファイルのみを対象"
        logger.info(log_message)
        if status_queue:
            status_queue.put(log_message)

    if exclude_patterns:
        log_message = f"除外パターン: {', '.join(exclude_patterns)}"
        logger.info(log_message)
        if status_queue:
            status_queue.put(log_message)

    # 指定されたフォルダパスとそのサブフォルダを検索
    for dirpath, dirnames, filenames in os.walk(target_folder_path):
        current_dir = Path(dirpath)

        # 状態更新
        if status_queue:
            status_queue.put(f"スキャン中: {current_dir}")

        # 除外パターンに一致するフォルダはスキップ (サブフォルダも含む)
        if match_exclude_pattern(current_dir, exclude_patterns):
            log_message = f"除外パターンに一致するフォルダをスキップ: {current_dir}"
            logger.info(log_message)
            if status_queue:
                status_queue.put(log_message)
            # このディレクトリ配下を無視するために dirnames をクリア
            dirnames[:] = []
            continue

        # フォルダ内のファイルを処理
        for file in filenames:
            file_path = current_dir / file
            full_path = str(file_path)

            # 除外パターンに一致する場合はスキップ
            if match_exclude_pattern(file_path, exclude_patterns):
                logger.debug(f"除外パターンに一致するファイルをスキップ: {file_path}")
                continue

            # 期間チェック
            if is_file_in_date_range(full_path, days):
                found_files.append(full_path)
                logger.debug(f"ファイル追加: {full_path}")
            else:
                logger.debug(f"期間外ファイルスキップ: {full_path}")

    log_message = f"検索完了: {len(found_files)} 個のファイルを発見"
    logger.info(log_message)
    if status_queue:
        status_queue.put(log_message)

    # ファイルが見つかった場合のみ、フォルダ情報を返す形式に合わせる (後方互換性のため)
    # ただし、この関数を使う側でファイルリストを直接使うように変更する方が良い
    if found_files:
        return [{"folder_path": str(target_folder_path), "files": found_files}]
    else:
        return []

def create_markdown_content(target_folder_path, found_folders_info, status_queue=None):
    """
    検索結果をMarkdown形式で出力 (指定フォルダパス基準)
    """
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    target_folder_name = Path(target_folder_path).name # フォルダ名を取得

    # 生成日時とフォルダ数はログにのみ記録
    log_message = f"生成日時: {timestamp}"
    logger.info(log_message)
    if status_queue:
        status_queue.put(log_message)

    # found_folders_info はリストだが、要素は1つのはず
    if not found_folders_info:
         md_content = f"# {target_folder_name} フォルダ内容集約\n\n"
         md_content += f"**指定されたフォルダパス**: `{target_folder_path}`\n\n"
         md_content += "**期間内のファイルは見つかりませんでした。**\n"
         return md_content

    # 最初の要素（唯一のはず）からファイルリストを取得
    files = found_folders_info[0].get("files", [])
    file_count = len(files)

    log_message = f"対象ファイル数: {file_count}"
    logger.info(log_message)
    if status_queue:
        status_queue.put(log_message)

    md_content = f"# {target_folder_name} フォルダ内容集約\n\n"
    md_content += f"**指定されたフォルダパス**: `{target_folder_path}`\n\n"
    md_content += f"**抽出されたファイル数**: {file_count}\n\n"

    if not files:
        md_content += "**期間内のファイルは見つかりませんでした。**\n"
        return md_content

    md_content += "## 目次\n\n"

    # ファイルを最終更新日時で降順ソート
    sorted_files = sorted(files, key=os.path.getmtime, reverse=True)

    # 目次作成 - ファイル名のみ表示 (フォルダ階層は複雑なので省略)
    for i, file_path in enumerate(sorted_files, 1):
        file_name = Path(file_path).name
        # アンカー用にファイル名を少し加工 (一意性を保証するのは難しいが簡易的に)
        file_anchor = f"file-{i}-{file_name.lower().replace(' ', '-').replace('.', '-')}"
        md_content += f"{i}. [{file_name}](#{file_anchor})\n"

    md_content += "\n---\n\n"

    # ファイルごとの詳細 - 新しいものから順に
    md_content += "## ファイル内容\n\n"
    for i, file_path in enumerate(sorted_files, 1):
        file_extension = Path(file_path).suffix.lower()
        file_name = Path(file_path).name
        file_mtime = datetime.datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M:%S")
        # アンカー名を目次と合わせる
        file_anchor = f"file-{i}-{file_name.lower().replace(' ', '-').replace('.', '-')}"

        if status_queue:
            status_queue.put(f"ファイル内容読込中: {file_path}")

        # Markdownの見出しにアンカーを追加
        md_content += f"### {i}. {file_name} <a name='{file_anchor}'></a>\n\n" # 見出しに番号追加
        md_content += f"**パス**: `{file_path}`\n"
        md_content += f"**最終更新日時**: {file_mtime}\n\n"

        # ファイル内容をコードブロックで表示
        file_content = read_file_content(file_path)

        # ファイル拡張子に応じた言語指定
        lang = ""
        if file_extension in ['.py', '.pyw']:
            lang = "python"
        elif file_extension in ['.md', '.markdown']:
            lang = "markdown"
        elif file_extension in ['.js', '.jsx']:
            lang = "javascript"
        elif file_extension in ['.html', '.htm']:
            lang = "html"
        elif file_extension in ['.css']:
            lang = "css"
        elif file_extension in ['.json']:
            lang = "json"
        elif file_extension in ['.xml']:
            lang = "xml"
        elif file_extension in ['.sh', '.bash']:
            lang = "bash"
        elif file_extension in ['.bat', '.cmd', '.ps1']:
            lang = "powershell"
        elif file_extension in ['.sql']:
            lang = "sql"
        elif file_extension == '.csv':
            lang = "csv"

        md_content += f"```{lang}\n{file_content}\n```\n\n"
        md_content += "---\n\n"

    if status_queue:
        status_queue.put("マークダウン生成完了")

    return md_content

# モダンなUI用のカスタムウィジェットクラス
class RoundedButton(tk.Canvas):
    """角丸のモダンなボタン"""
    def __init__(self, parent, text, command=None, width=100, height=36, 
                 bg_color="#2F4F4F", fg_color="white", font_family="Arial", 
                 font_size=10, font_weight="bold", corner_radius=10, hover_color="#3A6363", **kwargs):
        super().__init__(parent, width=width, height=height, 
                        bg=parent["bg"], bd=0, highlightthickness=0, **kwargs)
        
        self.bg_color = bg_color
        self.fg_color = fg_color
        self.hover_color = hover_color
        self.corner_radius = corner_radius
        self.current_color = bg_color
        self.command = command
        self.text = text
        
        # ボタンの描画
        self.button_bg = self.create_rounded_rect(0, 0, width, height, corner_radius, fill=bg_color)
        self.button_text = self.create_text(width//2, height//2, text=text, fill=fg_color,
                                           font=(font_family, font_size, font_weight))
        
        # イベントバインド
        self.bind("<Enter>", self.on_hover)
        self.bind("<Leave>", self.on_leave)
        self.bind("<ButtonPress-1>", self.on_press)
        self.bind("<ButtonRelease-1>", self.on_release)
    
    def create_rounded_rect(self, x1, y1, x2, y2, r, **kwargs):
        """角丸の長方形を描画"""
        points = [
            x1+r, y1,
            x2-r, y1,
            x2, y1,
            x2, y1+r,
            x2, y2-r,
            x2, y2,
            x2-r, y2,
            x1+r, y2,
            x1, y2,
            x1, y2-r,
            x1, y1+r,
            x1, y1
        ]
        return self.create_polygon(points, smooth=True, **kwargs)
    
    def on_hover(self, event):
        """ホバー状態の処理"""
        self.current_color = self.hover_color
        self.itemconfig(self.button_bg, fill=self.hover_color)
    
    def on_leave(self, event):
        """通常状態に戻る処理"""
        self.current_color = self.bg_color
        self.itemconfig(self.button_bg, fill=self.bg_color)
    
    def on_press(self, event):
        """ボタンプレス時の処理"""
        self.itemconfig(self.button_bg, fill=self.fg_color)
        self.itemconfig(self.button_text, fill=self.bg_color)
    
    def on_release(self, event):
        """ボタンリリース時の処理"""
        self.itemconfig(self.button_bg, fill=self.current_color)
        self.itemconfig(self.button_text, fill=self.fg_color)
        if self.command:
            self.command()

class ModernEntryFrame(tk.Frame):
    """モダンな入力欄を持つフレーム"""
    def __init__(self, parent, label_text, default_value="", width=250, **kwargs):
        super().__init__(parent, bg=parent["bg"], **kwargs)
        
        self.var = tk.StringVar(value=default_value)
        
        # ラベル
        self.label = tk.Label(self, text=label_text, bg=self["bg"], 
                            fg="#333", font=("Arial", 10, "bold"))
        self.label.pack(anchor="w", pady=(5, 2))
        
        # 入力フレーム（影のエフェクト用）
        self.entry_frame = tk.Frame(self, bg="#DDD", padx=1, pady=1)
        self.entry_frame.pack(fill="x", pady=(0, 5))
        
        # 実際の入力欄
        self.entry = tk.Entry(self.entry_frame, textvariable=self.var, 
                            font=("Arial", 10), bd=0, bg="white", fg="black")
        self.entry.pack(fill="x", ipady=8, padx=1, pady=1)
    
    def get(self):
        """入力された値を取得"""
        return self.var.get()
    
    def set(self, value):
        """値を設定"""
        self.var.set(value)

class ModernComboFrame(tk.Frame):
    """モダンなコンボボックスを持つフレーム"""
    def __init__(self, parent, label_text, values, default_index=0, width=250, **kwargs):
        super().__init__(parent, bg=parent["bg"], **kwargs)
        
        # 変数
        self.var = tk.StringVar()
        
        # ラベル
        self.label = tk.Label(self, text=label_text, bg=self["bg"], 
                            fg="#333", font=("Arial", 10, "bold"))
        self.label.pack(anchor="w", pady=(5, 2))
        
        # スタイル設定
        style = ttk.Style()
        style.configure("Modern.TCombobox", padding=8)
        
        # コンボボックス
        self.combo = ttk.Combobox(self, textvariable=self.var, values=values, 
                                width=width, style="Modern.TCombobox", state="readonly")
        # Mac対応：スタイルを明示的に設定
        style.configure("Modern.TCombobox", fieldbackground="white", foreground="black")
        self.combo.pack(fill="x", pady=(0, 5))
        
        # デフォルト値を設定
        if values and len(values) > default_index:
            self.combo.current(default_index)
    
    def get(self):
        """選択された値を取得"""
        return self.var.get()
    
    def set(self, value):
        """値を設定"""
        self.var.set(value)

class ModernCheckbutton(tk.Checkbutton):
    """モダンなチェックボタン"""
    def __init__(self, parent, text, **kwargs):
        self.var = tk.BooleanVar(value=False)
        super().__init__(parent, text=text, variable=self.var, 
                        bg=parent["bg"], fg="#333", font=("Arial", 10),
                        activebackground=parent["bg"], selectcolor="#EEE", **kwargs)
    
    def is_checked(self):
        """チェック状態を取得"""
        return self.var.get()
    
    def set(self, value):
        """チェック状態を設定"""
        self.var.set(value)

class ModernContextExtractorUI:
    """モダンなContext Extracter UI"""
    def __init__(self, root):
        self.root = root
        self.root.title("Context Extracter")
        
        # ウィンドウサイズと位置を設定
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        window_width = 700
        window_height = 800
        x_position = (screen_width - window_width) // 2
        y_position = (screen_height - window_height) // 2
        self.root.geometry(f"{window_width}x{window_height}+{x_position}+{y_position}")
        
        # BMSルートディレクトリを取得（Macでは正しく動作するように修正）
        self.bms_root_dir = get_bms_root()
        logger.info(f"BMSルートディレクトリ: {self.bms_root_dir}")
        
        # 背景色を設定
        self.bg_color = "#F5F5F7"
        self.accent_color = "#2F4F4F"  # チャコールグレー
        self.root.configure(bg=self.bg_color)
        
        # メインコンテナ
        self.main_container = tk.Frame(root, bg=self.bg_color, padx=20, pady=20)
        self.main_container.pack(fill=tk.BOTH, expand=True)
        
        # ヘッダーフレーム
        self.header_frame = tk.Frame(self.main_container, bg=self.bg_color)
        self.header_frame.pack(fill=tk.X, pady=(0, 20))
        
        # タイトル
        title_label = tk.Label(self.header_frame, text="Context Extracter", 
                             font=("Arial", 18, "bold"), bg=self.bg_color, fg=self.accent_color)
        title_label.pack(side=tk.LEFT)
        
        # 条件入力フレーム
        self.input_frame = tk.Frame(self.main_container, bg=self.bg_color, padx=10, pady=10)
        self.input_frame.pack(fill=tk.X)
        
        # 影付きの条件入力エリア（カード風）
        self.card_frame = tk.Frame(self.input_frame, bg="white", padx=15, pady=15,
                                 highlightbackground="#DDD", highlightthickness=1)
        self.card_frame.pack(fill=tk.X)
        
        # 基本入力項目
        self.root_entry = ModernEntryFrame(self.card_frame, "フォルダパス指定",
                                         default_value="", width=300)
        self.root_entry.pack(fill=tk.X, pady=5)
        
        # ルートディレクトリ参照ボタン
        browse_button = RoundedButton(self.card_frame, "参照",
                                    command=self.browse_root_dir, width=60, height=30,
                                    bg_color="#999", hover_color="#777")
        browse_button.pack(anchor="e", padx=5, pady=5)
        
        # 区切り線
        separator = tk.Frame(self.card_frame, height=1, bg="#EEE")
        separator.pack(fill=tk.X, pady=10)
        
        # 詳細設定ラベル
        detail_label = tk.Label(self.card_frame, text="詳細条件設定", 
                              font=("Arial", 10, "bold"), bg="white", fg="#333")
        detail_label.pack(anchor="w", pady=(5, 10))
        
        # 期間設定フレーム
        period_frame = tk.Frame(self.card_frame, bg="white")
        period_frame.pack(fill=tk.X, pady=5)
        
        # 期間設定ラベル
        period_label = tk.Label(period_frame, text="期間設定:", 
                              font=("Arial", 10), bg="white", fg="#333")
        period_label.pack(side=tk.LEFT, padx=(0, 10))
        
        # 日数入力欄
        self.days_var = tk.StringVar(value="14")
        self.days_entry = tk.Entry(period_frame, textvariable=self.days_var, 
                                 width=5, font=("Arial", 10), bd=1)
        self.days_entry.pack(side=tk.LEFT)
        
        # 日数ラベル
        days_label = tk.Label(period_frame, text="日以内", 
                            font=("Arial", 10), bg="white", fg="#333")
        days_label.pack(side=tk.LEFT, padx=(5, 20))
        
        # 全期間チェックボックス
        self.all_days_check = ModernCheckbutton(period_frame, "すべての期間")
        self.all_days_check.pack(side=tk.LEFT)
        
        # 除外パターンフレーム
        exclude_frame = tk.Frame(self.card_frame, bg="white")
        exclude_frame.pack(fill=tk.X, pady=10)
        
        # 除外パターンラベル
        exclude_label = tk.Label(exclude_frame, text="除外パターン:", 
                               font=("Arial", 10), bg="white", fg="#333")
        exclude_label.pack(side=tk.LEFT, padx=(0, 10))
        
        # 除外パターン入力欄
        self.exclude_var = tk.StringVar()
        self.exclude_entry = tk.Entry(exclude_frame, textvariable=self.exclude_var, 
                                    font=("Arial", 10), bd=1)
        self.exclude_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        # 除外パターン説明
        exclude_hint = tk.Label(exclude_frame, text="（スペース区切り）", 
                              font=("Arial", 8), bg="white", fg="#777")
        exclude_hint.pack(side=tk.RIGHT, padx=(5, 0))
        
        # 詳細ログチェックボックス
        self.verbose_check = ModernCheckbutton(self.card_frame, "詳細ログを出力する")
        self.verbose_check.pack(anchor="w", pady=5)
        
        # ボタンフレーム
        self.button_frame = tk.Frame(self.main_container, bg=self.bg_color, pady=15)
        self.button_frame.pack(fill=tk.X)
        
        # 抽出開始ボタン
        self.extract_button = RoundedButton(self.button_frame, "抽出開始", 
                                          command=self.run_extraction, width=120, height=40,
                                          bg_color=self.accent_color, hover_color="#3A6363",
                                          font_size=12)
        self.extract_button.pack(side=tk.LEFT, padx=(0, 10))
        
        # 結果フレーム
        self.result_frame = tk.Frame(self.main_container, bg=self.bg_color, pady=10)
        self.result_frame.pack(fill=tk.BOTH, expand=True)
        
        # 結果ヘッダーフレーム
        result_header = tk.Frame(self.result_frame, bg=self.bg_color)
        result_header.pack(fill=tk.X, pady=(0, 5))
        
        # 結果ラベル
        result_label = tk.Label(result_header, text="結果がここに表示されます...", 
                              font=("Arial", 10), bg=self.bg_color, fg="#333")
        result_label.pack(side=tk.LEFT)
        
        # 文字数カウント表示用
        self.extracted_count_var = tk.StringVar(value="抽出した文字数: 0")
        extracted_count_label = tk.Label(result_header, textvariable=self.extracted_count_var, 
                                   font=("Arial", 9), bg=self.bg_color, fg="#666")
        extracted_count_label.pack(side=tk.LEFT, padx=(20, 0))
        
        # コピーボタン（結果ヘッダーに追加）
        self.copy_button = RoundedButton(result_header, "コピー", 
                                       command=self.copy_to_clipboard, width=70, height=28,
                                       bg_color=self.accent_color, hover_color="#3A6363",
                                       font_size=9)
        self.copy_button.pack(side=tk.RIGHT, padx=5)
        
        # 180,000文字コピーボタン
        self.copy_180k_button = RoundedButton(result_header, "180K文字コピー", 
                                       command=self.copy_180k_to_clipboard, width=100, height=28,
                                       bg_color=self.accent_color, hover_color="#3A6363",
                                       font_size=9)
        self.copy_180k_button.pack(side=tk.RIGHT, padx=5)
        
        # 900,000文字コピーボタン
        self.copy_900k_button = RoundedButton(result_header, "900K文字コピー", 
                                       command=self.copy_900k_to_clipboard, width=100, height=28,
                                       bg_color=self.accent_color, hover_color="#3A6363",
                                       font_size=9)
        self.copy_900k_button.pack(side=tk.RIGHT, padx=5)
        
        # 結果テキストエリア（影付きカード風）
        self.result_container = tk.Frame(self.result_frame, bg="#DDD", padx=1, pady=1)
        self.result_container.pack(fill=tk.BOTH, expand=True)
        
        self.result_text = tk.Text(self.result_container, font=("Consolas", 10), 
                                  wrap=tk.WORD, bd=0, padx=10, pady=10,
                                  bg="white", fg="black", height=25)  # 高さを増やしました
        self.result_text.pack(fill=tk.BOTH, expand=True, side=tk.LEFT)
        
        # スクロールバー
        scrollbar = tk.Scrollbar(self.result_container)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.result_text.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.result_text.yview)
        
        # テキスト文字数カウント用フレーム
        self.count_frame = tk.Frame(self.result_frame, bg=self.bg_color)
        self.count_frame.pack(fill=tk.X, pady=(5, 0))
        
        # 文字数カウントラベル
        self.char_count_var = tk.StringVar(value="文字数: 0")
        char_count_label = tk.Label(self.count_frame, textvariable=self.char_count_var, 
                                   font=("Arial", 9), bg=self.bg_color, fg="#666")
        char_count_label.pack(side=tk.LEFT)
        
        # テキスト変更イベントをバインド
        self.result_text.bind("<<Modified>>", self.update_char_count)
        
        # 操作ボタンフレーム
        self.action_frame = tk.Frame(self.main_container, bg=self.bg_color, pady=10)
        self.action_frame.pack(fill=tk.X)
        
        # 保存ボタン
        self.save_button = RoundedButton(self.action_frame, "保存", 
                                       command=self.save_to_file, width=100, height=36,
                                       bg_color=self.accent_color, hover_color="#3A6363")
        self.save_button.pack(side=tk.RIGHT)
        
        # ステータスバー
        self.status_var = tk.StringVar(value="準備完了")
        self.status_bar = tk.Label(self.root, textvariable=self.status_var, 
                                  bg="#E5E5E7", fg="#666", font=("Arial", 9),
                                  bd=1, relief=tk.SUNKEN, anchor=tk.W, padx=10, pady=2)
        self.status_bar.pack(side=tk.BOTTOM, fill=tk.X)
        
        # キューの初期化
        self.status_queue = queue.Queue()
        self.result_queue = queue.Queue()
        
        # キューチェックのスケジュール
        self.root.after(100, self.check_queue)
        
        # 初期フォーカスを設定
        self.root_entry.entry.focus_set() # フォルダパス入力にフォーカス
    
    def browse_root_dir(self):
        """ルートディレクトリを選択するダイアログを表示"""
        directory = filedialog.askdirectory(initialdir=self.root_entry.get())
        if directory:
            self.root_entry.set(directory)
    
    def check_queue(self):
        """キューに溜まったメッセージを処理"""
        try:
            # ステータスキューを処理
            while not self.status_queue.empty():
                message = self.status_queue.get(0)
                self.status_var.set(message)
                
            # 結果キューを処理
            while not self.result_queue.empty():
                result = self.result_queue.get(0)
                self.set_result(result)
        except Exception as e:
            logger.error(f"キュー処理エラー: {e}")
        finally:
            # 再スケジュール
            self.root.after(100, self.check_queue)
    
    def set_result(self, text):
        """結果テキストに文字を設定"""
        self.result_text.delete(1.0, tk.END)
        self.result_text.insert(tk.END, text)
        # テキスト色を明示的に設定
        self.result_text.config(fg="black")
        # 文字数更新
        self.update_char_count()
    
    def update_char_count(self, event=None):
        """テキストエリアの文字数をカウントして表示"""
        try:
            # 文字数を取得
            text = self.result_text.get("1.0", tk.END)
            char_count = len(text) - 1  # ENDマーカーには\nが含まれるので-1
            
            # 文字数と行数を表示
            lines = text.count('\n')
            self.char_count_var.set(f"文字数: {char_count} | 行数: {lines}")
            
            # 抽出した文字数も更新
            self.extracted_count_var.set(f"抽出した文字数: {char_count}")
            
            # Modifiedフラグをリセット
            self.result_text.edit_modified(False)
        except Exception as e:
            logger.error(f"文字数カウントエラー: {e}")
    
    def copy_to_clipboard(self):
        """結果をクリップボードにコピー"""
        try:
            result = self.result_text.get(1.0, tk.END)
            if result.strip():
                self.root.clipboard_clear()
                self.root.clipboard_append(result)
                self.status_var.set("結果をクリップボードにコピーしました")
            else:
                messagebox.showinfo("情報", "コピーする結果がありません")
        except Exception as e:
            messagebox.showerror("エラー", f"クリップボードへのコピーに失敗しました: {e}")
            
    def copy_180k_to_clipboard(self):
        """結果の最初の180,000文字をクリップボードにコピー"""
        try:
            result = self.result_text.get(1.0, tk.END)
            if result.strip():
                # 最初の180,000文字を取得
                text_to_copy = result[:180000]
                self.root.clipboard_clear()
                self.root.clipboard_append(text_to_copy)
                self.status_var.set("最初の180,000文字をクリップボードにコピーしました")
            else:
                messagebox.showinfo("情報", "コピーする結果がありません")
        except Exception as e:
            messagebox.showerror("エラー", f"クリップボードへのコピーに失敗しました: {e}")
    
    def copy_900k_to_clipboard(self):
        """結果の最初の900,000文字をクリップボードにコピー"""
        try:
            result = self.result_text.get(1.0, tk.END)
            if result.strip():
                # 最初の900,000文字を取得
                text_to_copy = result[:900000]
                self.root.clipboard_clear()
                self.root.clipboard_append(text_to_copy)
                self.status_var.set("最初の900,000文字をクリップボードにコピーしました")
            else:
                messagebox.showinfo("情報", "コピーする結果がありません")
        except Exception as e:
            messagebox.showerror("エラー", f"クリップボードへのコピーに失敗しました: {e}")
    
    def save_to_file(self):
        """結果をファイルに保存"""
        try:
            # フォルダパスを取得
            target_folder_path_str = self.root_entry.get().strip()
            if not target_folder_path_str:
                 messagebox.showerror("エラー", "フォルダパスを指定してください。")
                 return

            target_folder_path = Path(target_folder_path_str)
            folder_name = target_folder_path.name # パスの最後の部分をフォルダ名として使用
            if not folder_name: # ルートディレクトリなどが指定された場合
                folder_name = "root_context"

            # output_dir の存在を確認 (グローバル変数を使用)
            global output_dir
            if not output_dir.exists():
                logger.warning(f"出力ディレクトリが見つかりません: {output_dir}. 作成を試みます。")
                try:
                    output_dir.mkdir(parents=True, exist_ok=True)
                except Exception as mkdir_e:
                    messagebox.showerror("エラー", f"出力ディレクトリの作成に失敗しました: {output_dir}\nエラー: {mkdir_e}")
                    return

            # 現在時刻を含めたファイル名
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            # ファイル名に使えない文字を置換
            safe_folder_name = re.sub(r'[\\/*?:"<>|]', '_', folder_name)
            default_filename = f"{safe_folder_name}_{timestamp}_contents.md" # フォルダ名を使用

            filepath = filedialog.asksaveasfilename(
                initialdir=str(output_dir), # 文字列に変換して渡す
                initialfile=default_filename,
                defaultextension=".md",
                filetypes=[("Markdown files", "*.md"), ("All files", "*.*")]
            )

            if not filepath:
                self.status_var.set("ファイル保存がキャンセルされました")
                return

            # ファイルパスをPathオブジェクトに変換
            filepath_obj = Path(filepath)

            with open(filepath_obj, 'w', encoding='utf-8') as f:
                content = self.result_text.get(1.0, tk.END)
                # Tkinter Textウィジェットは末尾に改行を追加する場合があるので、不要なら削除
                if content.endswith('\n'):
                    content = content[:-1]
                f.write(content)

            self.status_var.set(f"結果をファイルに保存しました: {filepath_obj}")
            logger.info(f"結果をファイルに保存しました: {filepath_obj}")

        except Exception as e:
            error_msg = f"ファイルの保存に失敗しました: {e}"
            logger.error(error_msg, exc_info=True) # 詳細なエラー情報をログに出力
            messagebox.showerror("エラー", error_msg)
            self.status_var.set("ファイルの保存中にエラーが発生しました")
    
    def run_extraction(self):
        """抽出処理を実行"""
        # 入力値の検証
        target_folder_path = self.root_entry.get().strip() # フォルダパスを取得
        if not target_folder_path:
            messagebox.showerror("エラー", "フォルダパスを指定してください") # フォルダパス必須に
            return
            
        # 実行中であることをUIに表示
        self.result_text.delete(1.0, tk.END)
        self.result_text.insert(tk.END, "抽出処理を実行中です...\nしばらくお待ちください。")
        self.result_text.config(fg="black")
        self.root.update()

        # 指定されたパスが存在するか、ディレクトリかを確認
        target_path_obj = Path(target_folder_path)
        if not target_path_obj.exists():
            messagebox.showerror("エラー", f"指定されたパスが見つかりません:\n{target_folder_path}")
            return
        if not target_path_obj.is_dir():
             messagebox.showerror("エラー", f"指定されたパスはフォルダではありません:\n{target_folder_path}")
             return

        # 除外パターンを分割
        exclude_patterns = [p.strip() for p in self.exclude_var.get().split() if p.strip()]

        # 日数の処理
        if self.all_days_check.is_checked():
            days = None
        else:
            try:
                days = int(self.days_var.get())
                if days <= 0:
                    messagebox.showerror("エラー", "日数は正の整数を指定してください")
                    return
            except ValueError:
                messagebox.showerror("エラー", "日数に有効な整数を入力してください")
                return
        
        # 詳細ログの設定
        if self.verbose_check.is_checked():
            logger.setLevel(logging.DEBUG)
            for handler in logger.handlers:
                if isinstance(handler, logging.FileHandler):
                    handler.setLevel(logging.DEBUG)
        else:
            logger.setLevel(logging.INFO)
            for handler in logger.handlers:
                if isinstance(handler, logging.FileHandler):
                    handler.setLevel(logging.INFO)
        
        # ステータス更新
        self.status_var.set("抽出処理を開始します...")
        
        # 結果エリアをクリア
        self.result_text.delete(1.0, tk.END)
        
        # スレッドで処理を実行
        threading.Thread(target=self.extraction_thread, args=(
            target_folder_path, days, exclude_patterns)).start() # 引数変更
    
    def extraction_thread(self, target_folder_path, days, exclude_patterns): # 引数変更
        """別スレッドで実行する抽出処理"""
        try:
            # フォルダとファイルを検索
            root_path = Path(target_folder_path) # 名前変更
            self.status_queue.put(f"検索中: {target_folder_path}")
            # 関数呼び出し変更 & 引数変更
            found_folders_info = find_folders_and_files(root_path, days, exclude_patterns, self.status_queue)

            # Markdown ファイルの作成
            self.status_queue.put(f"マークダウン生成中...")
            # 関数呼び出し変更 & 引数変更
            md_content = create_markdown_content(target_folder_path, found_folders_info, self.status_queue)

            # 結果を表示
            self.result_queue.put(md_content)

            # 完了メッセージと結果の文字数をログに記録
            char_count = len(md_content)
            completion_message = f"抽出処理が完了しました - 合計文字数: {char_count}"
            logger.info(completion_message)
            self.status_queue.put(completion_message)

        except Exception as e:
            error_message = f"エラーが発生しました: {e}"
            logger.error(error_message, exc_info=True)
            self.status_queue.put(error_message)

def main():
    """アプリケーションのメインエントリポイント"""
    root = tk.Tk()
    app = ModernContextExtractorUI(root)
    root.mainloop()

if __name__ == "__main__":
    main()
