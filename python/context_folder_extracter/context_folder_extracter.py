#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Context Folder Extracter (Command Line)

指定したフォルダパスとそのサブフォルダ内のファイル内容を抽出。
すべてのファイル内容を1つのMarkdownファイルに集約する。

使用例:
        python context_folder_extracter.py <target_folder_path> [options]

引数:
        target_folder_path: 内容を抽出するフォルダのパス (必須)

オプション:
        --days, -d DAYS: 過去何日以内のファイルを対象にするか (デフォルト: 14)
        --all, -a:     すべての期間のファイルを対象にする (--daysより優先)
        --exclude, -e [PATTERN ...]: 除外するフォルダやファイルのパターン
        --verbose, -v: 詳細なログを出力する

出力:
        output ディレクトリ内に、<フォルダ名>_<タイムスタンプ>_contents.md というファイル名で出力
"""

import os
import sys
import platform
import datetime
import time
import argparse
from pathlib import Path
import logging
import re
import csv # csv モジュールをインポート

# 必要なライブラリをインポート (UI版と同様)
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    print("警告: 'openpyxl' ライブラリが見つかりません。Excel (.xlsx, .xlsm) ファイルの処理はスキップされます。`pip install openpyxl` でインストールしてください。", file=sys.stderr)

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("警告: 'python-pptx' ライブラリが見つかりません。PowerPoint (.pptx) ファイルの処理はスキップされます。`pip install python-pptx` でインストールしてください。", file=sys.stderr)

try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False
    print("警告: 'PyPDF2' ライブラリが見つかりません。PDF (.pdf) ファイルの処理はスキップされます。`pip install PyPDF2` でインストールしてください。", file=sys.stderr)

# pywin32 のインポート (Windowsのみ)
import platform # platform をインポート (重複インポート防止のため確認)
HAS_PYWIN32 = False
if platform.system() == "Windows":
    try:
        import win32com.client
        import pythoncom # COM初期化に必要
        HAS_PYWIN32 = True
    except ImportError:
        print("警告: 'pywin32' ライブラリが見つかりません。.ppt ファイルの処理はスキップされます。Windows環境の場合 `pip install pywin32` でインストールしてください。", file=sys.stderr)
else:
    # Windows以外の場合も警告を表示
    print("情報: Windows以外の環境では .ppt ファイルの処理はサポートされません。", file=sys.stderr)

# output ディレクトリの存在確認または作成
# __file__ から絶対パスを取得し、その親ディレクトリを基準にする
try:
    # sys.argv[0]を使用してスクリプトの絶対パスを取得（MacとWindowsの両方で動作）
    script_dir = Path(os.path.abspath(sys.argv[0])).parent
    output_dir = script_dir / "output"
    output_dir.mkdir(exist_ok=True)
    log_file_path = output_dir / "context_folder_extracter.log"
except Exception as e:
    # outputディレクトリ作成やログファイルパス設定に失敗した場合
    # スクリプトと同じディレクトリにログを出力するなどのフォールバック処理
    script_dir = Path(__file__).resolve().parent
    log_file_path = script_dir / "context_folder_extracter.fallback.log"
    print(f"警告: outputディレクトリの準備に失敗しました。ログは {log_file_path} に出力されます。エラー: {e}", file=sys.stderr)
    # output_dir が未定義の場合に備え、スクリプトディレクトリを割り当てておく
    if 'output_dir' not in locals():
        output_dir = script_dir

# ロガーの設定
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(log_file_path, mode="w", encoding="utf-8"), # 上書きモードに設定
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

logger.info(f"スクリプトディレクトリ: {script_dir}")
logger.info(f"出力ディレクトリ: {output_dir}")
logger.info(f"ログファイルパス: {log_file_path}")

def get_bms_root():
    """
    BMSのルートディレクトリを取得する。環境に合わせて調整し、存在確認を行う。
    """
    system = platform.system()
    bms_root = None
    if system == "Windows":
        # Windowsの場合は D:/BMS を想定
        bms_root = Path("D:/BMS")
    elif system == "Darwin": # macOS の正式名は Darwin
        # macOSの場合はホームディレクトリ下の Desktop/BMS を想定
        home_dir = Path.home() # ホームディレクトリを取得
        bms_root = home_dir / "Desktop" / "BMS"
        logger.info(f"macOS用BMSルートディレクトリを設定: {bms_root}")
    else: # その他のLinuxなど
        # 必要であれば他のOS用のパスも定義
        logger.warning(f"未対応のOSです: {system}. デフォルトパスは設定されません。")
        return None # PathオブジェクトではなくNoneを返す

    # パスが存在するか確認
    if bms_root and not bms_root.exists():
        logger.warning(f"デフォルトのBMSルートディレクトリが見つかりません: {bms_root}")
        # 存在しない場合でもパス自体は返す

    return bms_root # Pathオブジェクトを返す

def is_binary_file(file_path):
    """
    ファイルがバイナリかテキストかを判定
    ただし、特定の拡張子（xlsx, xlsm, pptx, ppt, pdf, csv）は内容抽出を試みるため除外
    """
    file_path_obj = Path(file_path)
    ext = file_path_obj.suffix.lower()

    # 特定の拡張子はバイナリ判定を行わず、専用の読み込み処理に任せる
    if ext in ['.xlsx', '.xlsm', '.pptx', '.ppt', '.pdf', '.csv']:
        return False # バイナリではない扱いにする（実際はバイナリだが、専用処理があるため）

    textchars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7f})
    try:
        with open(file_path, 'rb') as f:
            is_binary = bool(f.read(1024).translate(None, textchars))
        return is_binary
    except Exception as e:
        logger.warning(f"ファイル {file_path} の判定中にエラー発生: {e}")
        return True  # エラーが発生した場合もバイナリとして扱う

def is_file_in_date_range(file_path, days=None):
    """
    ファイルが指定された日数以内かどうか確認
    daysがNoneの場合は、すべてのファイルが対象
    """
    if days is None:
        return True
    
    try:
        # ファイルの最終更新日時を取得
        file_mtime = os.path.getmtime(file_path)
        now = time.time()
        
        # 日数を秒に変換
        days_in_seconds = days * 24 * 60 * 60
        
        # 指定期間内かチェック
        return (now - file_mtime) <= days_in_seconds
    except Exception as e:
        logger.warning(f"ファイル日付の確認中にエラー発生: {file_path} - {e}")
        return True  # エラーが発生した場合は含める

def match_exclude_pattern(path, exclude_patterns):
    """
    指定された除外パターンにパスが一致するかどうか確認
    いずれかのパターンに一致すればTrueを返す
    """
    if not exclude_patterns:
        return False
    
    path_lower = str(path).lower()
    for pattern in exclude_patterns:
        if pattern.lower() in path_lower:
            logger.debug(f"除外パターン '{pattern}' がパス '{path}' に一致")
            return True
    return False

def read_excel_content(file_path):
    """Excel (.xlsx, .xlsm) ファイルの内容を抽出"""
    if not HAS_OPENPYXL:
        return f"[ライブラリ不足のためスキップ (openpyxl): {file_path}]"
    try:
        # .xlsm も openpyxl で基本的なテキストは読める想定
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        content = []
        for sheet_name in workbook.sheetnames:
            content.append(f"--- シート: {sheet_name} ---")
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                row_content = []
                for cell in row:
                    if cell.value is not None:
                        row_content.append(str(cell.value))
                if row_content:
                    content.append("\t".join(row_content))
        return "\n".join(content)
    except Exception as e:
        logger.warning(f"Excelファイル {file_path} の読込中にエラー発生: {e}")
        # xlsm特有のエラー考慮 (例: 拡張機能関連)
        if ".xlsm" in str(file_path).lower() and "macros" in str(e).lower():
             return f"[Excel(マクロ有効)ファイル読込エラー(マクロ/拡張機能関連の可能性): {file_path}]"
        return f"[Excelファイル読込エラー: {file_path}]"

def read_pptx_content(file_path):
    """PowerPoint (.pptx) ファイルの内容を抽出"""
    if not HAS_PPTX:
        return f"[ライブラリ不足のためスキップ (python-pptx): {file_path}]"
    try:
        prs = Presentation(file_path)
        content = []
        for i, slide in enumerate(prs.slides):
            content.append(f"--- スライド {i+1} ---")
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_texts.append(shape.text)
            # ノートも取得する場合
            if slide.has_notes_slide:
                 notes_slide = slide.notes_slide
                 notes_text = notes_slide.notes_text_frame.text
                 if notes_text:
                     content.append(f"[ノート]\n{notes_text}")

            if slide_texts:
                content.append("\n".join(slide_texts))

        return "\n".join(content)
    except Exception as e:
        logger.warning(f"PowerPointファイル {file_path} の読込中にエラー発生: {e}")
        return f"[PowerPointファイル読込エラー: {file_path}]"

def read_pdf_content(file_path):
    """PDF (.pdf) ファイルの内容を抽出"""
    if not HAS_PYPDF2:
        return f"[ライブラリ不足のためスキップ (PyPDF2): {file_path}]"
    try:
        content = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            num_pages = len(reader.pages)
            for i in range(num_pages):
                page = reader.pages[i]
                page_text = page.extract_text()
                if page_text: # テキストが存在する場合のみ追加
                     content.append(f"--- ページ {i+1} ---")
                     content.append(page_text)
        # PyPDF2がテキストを抽出できなかった場合（画像PDFなど）
        if not content:
             return f"[PDFからテキスト抽出不可（画像PDF等の可能性）: {file_path}]"

        return "\n".join(content)
    except Exception as e:
        logger.warning(f"PDFファイル {file_path} の読込中にエラー発生: {e}")
        # PyPDF2は暗号化されたPDFなどでエラーを出すことがある
        if "encrypted" in str(e).lower():
            return f"[暗号化されたPDFのため読込不可: {file_path}]"
        return f"[PDFファイル読込エラー: {file_path}]"

def read_ppt_content(file_path):
    """PowerPoint (.ppt) ファイルの内容を抽出 (Windowsのみ, 要PowerPointインストール)"""
    if platform.system() != "Windows":
        return f"[Windows環境ではないためスキップ (.ppt): {file_path}]"
    if not HAS_PYWIN32:
        return f"[ライブラリ不足のためスキップ (pywin32): {file_path}]"

    powerpoint = None
    presentation = None
    abs_path = str(Path(file_path).resolve()) # COMに渡すため絶対パスに変換

    try:
        # COMライブラリを初期化 (マルチスレッド環境で必要だがCLIでも安全のため)
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # PowerPointを非表示で実行
        powerpoint.Visible = False
        # ファイルを開く (読み取り専用、タイトル表示なし、ウィンドウ非表示)
        presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True, Untitled=False, WithWindow=False)

        content = []
        for i, slide in enumerate(presentation.Slides):
            content.append(f"--- スライド {i+1} ---")
            slide_texts = []
            try:
                for shape in slide.Shapes:
                    # テキストフレームを持つシェイプからテキストを抽出
                    if shape.HasTextFrame:
                        if shape.TextFrame.HasText:
                            slide_texts.append(shape.TextFrame.TextRange.Text)
            except Exception as shape_e:
                 logger.warning(f".pptファイル {file_path} スライド {i+1} のシェイプ処理中にエラー: {shape_e}")
                 slide_texts.append(f"[シェイプ読込エラー: {shape_e}]")

            # ノートページのテキストを取得
            try:
                if slide.HasNotesPage:
                    # ノートは通常 Placeholders(2) にある
                    if len(slide.NotesPage.Shapes.Placeholders) >= 2:
                        notes_shape = slide.NotesPage.Shapes.Placeholders[2]
                        if notes_shape.HasTextFrame:
                            if notes_shape.TextFrame.HasText:
                                notes_text = notes_shape.TextFrame.TextRange.Text
                                if notes_text.strip(): # 空でないノートのみ追加
                                    content.append(f"[ノート]\n{notes_text}")
                    else:
                         # 稀にプレースホルダー構成が異なる場合、他のテキストを持つシェイプを探す
                         for notes_shape in slide.NotesPage.Shapes:
                              if notes_shape.HasTextFrame and notes_shape.TextFrame.HasText:
                                   notes_text = notes_shape.TextFrame.TextRange.Text
                                   if notes_text.strip():
                                        content.append(f"[ノート(代替)]\n{notes_text}")
                                        break # 最初に見つかったものを採用

            except Exception as note_e:
                 logger.warning(f".pptファイル {file_path} スライド {i+1} のノート処理中にエラー: {note_e}")
                 content.append(f"[ノート読込エラー: {note_e}]")


            if slide_texts:
                content.append("\n".join(slide_texts))

        return "\n".join(content)
    except pythoncom.com_error as e:
         # COMエラーの詳細を出力
         logger.error(f".pptファイル {file_path} の読込中にCOMエラー発生: {e}")
         hresult = getattr(e, 'hresult', 'N/A')
         desc = getattr(e, 'excepinfo', ('','','',''))[2] if getattr(e, 'excepinfo', None) else "詳細不明"
         # PowerPointがインストールされていない場合などの特定のエラーを判別
         if hresult == -2147221005: # REGDB_E_CLASSNOTREG (クラスが登録されていない)
             return f"[PowerPointがインストールされていないか、COM登録に問題があります: {file_path}]"
         return f"[.pptファイル読込COMエラー (HRESULT: {hresult}, Desc: {desc}): {file_path}]"
    except Exception as e:
        logger.warning(f".pptファイル {file_path} の読込中に予期せぬエラー発生: {e}")
        return f"[.pptファイル読込エラー: {file_path}]"
    finally:
        try:
            if presentation:
                presentation.Close()
            if powerpoint:
                powerpoint.Quit()
        except Exception as close_e:
             logger.warning(f"PowerPoint終了処理中にエラー: {close_e}")
        # COMライブラリを解放
        pythoncom.CoUninitialize()

def read_csv_content(file_path):
    """CSV (.csv) ファイルの内容を抽出 (タブ区切りも考慮し、エンコーディングを試行)"""
    # 日本語でよく使われるエンコーディングを優先的に試す
    encodings_to_try = ['utf-8', 'shift-jis', 'cp932', 'euc-jp', 'iso2022-jp', 'utf-16']
    content = f"[CSVファイル: {file_path}]\n"
    detected_encoding = None
    detected_delimiter = None

    for encoding in encodings_to_try:
        try:
            with open(file_path, 'r', encoding=encoding, newline='') as f:
                # 最初の数行を読み込んで区切り文字を推定
                sample = "".join(f.readline() for _ in range(10)) # 判定のため少し多めに読む
                if not sample:
                    return f"[空のCSVファイル: {file_path}]"

                sniffer = csv.Sniffer()
                try:
                    dialect = sniffer.sniff(sample, delimiters=',\t;|') # 一般的な区切り文字候補
                    detected_delimiter = dialect.delimiter
                except csv.Error:
                    # 自動判定失敗時はカンマをデフォルトとするが、タブの可能性も考慮
                    if '\t' in sample and ',' not in sample:
                         detected_delimiter = '\t'
                    else:
                         detected_delimiter = ',' # デフォルトはカンマ

                f.seek(0) # ファイルポインタを先頭に戻す
                # 区切り文字を確定してリーダーを作成
                reader = csv.reader(f, delimiter=detected_delimiter)
                csv_data = []
                for row in reader:
                    # 各セルを文字列に変換し、空文字はそのまま保持
                    csv_data.append(detected_delimiter.join(map(str, row)))
                detected_encoding = encoding
                # 抽出内容の前にエンコーディングと区切り文字情報を追加
                content += f"--- (エンコーディング: {detected_encoding}, 区切り文字: '{detected_delimiter}') ---\n"
                content += "\n".join(csv_data)
                return content # 読み込めたらループを抜ける
        except UnicodeDecodeError:
            continue # 次のエンコーディングを試す
        except Exception as e:
            logger.warning(f"CSVファイル {file_path} (エンコーディング: {encoding}) の読込中にエラー発生: {e}")
            # 一度エラーが出ても他のエンコーディングで成功する可能性があるのでループは継続

    # すべてのエンコーディングで失敗した場合
    logger.warning(f"CSVファイル {file_path} の適切なエンコーディングが見つかりませんでした。試行したエンコーディング: {encodings_to_try}")
    return f"[CSVファイル読込エラー(エンコーディング不明またはファイル破損): {file_path}]"

def read_file_content(file_path):
    """
    ファイルの内容を読む。エラーで対応
    バイナリファイルの場合は、[バイナリファイル]と表示
    Excel, PowerPoint, PDF, CSV は専用関数で内容を抽出
    """
    file_path_obj = Path(file_path)
    ext = file_path_obj.suffix.lower()

    # 拡張子に応じて処理を分岐
    if ext == '.xlsx' or ext == '.xlsm': # .xlsm を .xlsx と同じ扱いに
        return read_excel_content(file_path)
    elif ext == '.pptx':
        return read_pptx_content(file_path)
    elif ext == '.ppt': # .ppt を追加
        return read_ppt_content(file_path)
    elif ext == '.pdf':
        return read_pdf_content(file_path)
    elif ext == '.csv': # .csv を追加
        return read_csv_content(file_path)
    # 既存のバイナリ判定とテキスト読み込み
    elif is_binary_file(file_path): # is_binary_fileがTrueを返すのは上記拡張子以外
        return f"[バイナリファイル: {file_path}]"
    else: # テキストファイルの場合
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='shift-jis') as f:
                    return f.read()
            except Exception as e:
                logger.warning(f"ファイル {file_path} の読込中にエラー発生 (Shift-JIS試行後): {e}")
                return f"[ファイル読込エラー(Shift-JIS試行後): {file_path}]"
        except Exception as e:
            logger.warning(f"ファイル {file_path} の読込中にエラー発生: {e}")
            return f"[ファイル読込エラー: {file_path}]"

def find_folders_and_files(target_folder_path, days=None, exclude_patterns=None):
    """
    指定されたフォルダパスとそのサブフォルダ内のファイルを検索
    days: 指定された日数以内のファイルのみを含める（Noneの場合はすべてのファイル）
    exclude_patterns: 除外するパターンのリスト
    """
    found_files = [] # ファイルパスのリスト
    
    # target_folder_pathがPathオブジェクトの場合は文字列に変換
    if isinstance(target_folder_path, Path):
        target_folder_path = str(target_folder_path)

    log_message = f"検索開始: フォルダパス '{target_folder_path}' 内を検索"
    logger.info(log_message)

    if days is not None:
        log_message = f"期間制限: 過去{days}日以内に更新されたファイルのみを対象"
        logger.info(log_message)

    if exclude_patterns:
        log_message = f"除外パターン: {', '.join(exclude_patterns)}"
        logger.info(log_message)

    # 指定されたフォルダパスとそのサブフォルダを検索
    for dirpath, dirnames, filenames in os.walk(target_folder_path):
        current_dir = Path(dirpath)

        # 除外パターンに一致するフォルダはスキップ (サブフォルダも含む)
        if match_exclude_pattern(current_dir, exclude_patterns):
            log_message = f"除外パターンに一致するフォルダをスキップ: {current_dir}"
            logger.info(log_message)
            # このディレクトリ配下を無視するために dirnames をクリア
            dirnames[:] = []
            continue

        # フォルダ内のファイルを処理
        for file in filenames:
            file_path = current_dir / file
            full_path = str(file_path)

            # 除外パターンに一致する場合はスキップ
            if match_exclude_pattern(file_path, exclude_patterns):
                logger.debug(f"除外パターンに一致するファイルをスキップ: {file_path}")
                continue

            # 期間チェック
            if is_file_in_date_range(full_path, days):
                found_files.append(full_path)
                logger.debug(f"ファイル追加: {full_path}")
            else:
                logger.debug(f"期間外ファイルスキップ: {full_path}")

    log_message = f"検索完了: {len(found_files)} 個のファイルを発見"
    logger.info(log_message)

    # ファイルが見つかった場合のみ、フォルダ情報を返す形式に合わせる
    if found_files:
        return [{"folder_path": str(target_folder_path), "files": found_files}]
    else:
        return []

def create_markdown_content(target_folder_path, found_folders_info): # 引数名変更
    """
    検索結果をMarkdown形式で出力 (指定フォルダパス基準)
    """
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    target_folder_name = Path(target_folder_path).name # フォルダ名を取得

    # 生成日時とファイル数はログにのみ記録
    logger.info(f"生成日時: {timestamp}")

    # found_folders_info はリストだが、要素は1つのはず
    if not found_folders_info:
         md_content = f"# {target_folder_name} フォルダ内容集約\n\n"
         md_content += f"**指定されたフォルダパス**: `{target_folder_path}`\n\n"
         md_content += "**期間内のファイルは見つかりませんでした。**\n"
         return md_content

    # 最初の要素（唯一のはず）からファイルリストを取得
    files = found_folders_info[0].get("files", [])
    file_count = len(files)

    logger.info(f"対象ファイル数: {file_count}")

    md_content = f"# {target_folder_name} フォルダ内容集約\n\n"
    md_content += f"**指定されたフォルダパス**: `{target_folder_path}`\n\n"
    md_content += f"**抽出されたファイル数**: {file_count}\n\n"

    if not files:
        md_content += "**期間内のファイルは見つかりませんでした。**\n"
        return md_content

    md_content += "## 目次\n\n"

    # ファイルを最終更新日時で降順ソート
    sorted_files = sorted(files, key=os.path.getmtime, reverse=True)

    # 目次作成 - ファイル名のみ表示
    for i, file_path in enumerate(sorted_files, 1):
        file_name = Path(file_path).name
        file_anchor = f"file-{i}-{file_name.lower().replace(' ', '-').replace('.', '-')}"
        md_content += f"{i}. [{file_name}](#{file_anchor})\n"

    md_content += "\n---\n\n"

    # ファイルごとの詳細 - 新しいものから順に
    md_content += "## ファイル内容\n\n"
    for i, file_path in enumerate(sorted_files, 1):
        file_extension = Path(file_path).suffix.lower()
        file_name = Path(file_path).name
        file_mtime = datetime.datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M:%S")
        file_anchor = f"file-{i}-{file_name.lower().replace(' ', '-').replace('.', '-')}"

        # Markdownの見出しにアンカーを追加
        md_content += f"### {i}. {file_name} <a name='{file_anchor}'></a>\n\n"
        md_content += f"**パス**: `{file_path}`\n"
        md_content += f"**最終更新日時**: {file_mtime}\n\n"

        # ファイル内容をコードブロックで表示
        file_content = read_file_content(file_path)

        # ファイル拡張子に応じた言語指定 (主要なもののみ抜粋)
        lang = ""
        if file_extension in ['.py', '.pyw']:
            lang = "python"
        elif file_extension in ['.md', '.markdown']:
            lang = "markdown"
        elif file_extension in ['.js', '.jsx']:
            lang = "javascript"
        elif file_extension in ['.html', '.htm']:
            lang = "html"
        elif file_extension in ['.css']:
            lang = "css"
        elif file_extension in ['.json']:
            lang = "json"
        elif file_extension in ['.xml']:
            lang = "xml"
        elif file_extension in ['.sh', '.bash']:
            lang = "bash"
        elif file_extension in ['.bat', '.cmd', '.ps1']:
            lang = "powershell"
        elif file_extension in ['.sql']:
            lang = "sql"
        elif file_extension == '.csv': # CSV を追加
            lang = "csv"
        # .ppt, .pptx, .xlsx, .xlsm, .pdf は特定の言語指定なし (プレーンテキスト扱い)

        md_content += f"```{lang}\n{file_content}\n```\n\n"
        md_content += "---\n\n"

    return md_content

def main():
    """
    メイン処理
    """
    parser = argparse.ArgumentParser(description='指定したフォルダパス内のファイル内容を集約するツール')
    parser.add_argument('target_folder_path', help='内容を抽出するフォルダのパス') # 必須の位置引数に変更
    # parser.add_argument('root_dir', nargs='?', default=None, help='...') # 削除
    parser.add_argument('--verbose', '-v', action='store_true', help='詳細なログを出力する')
    parser.add_argument('--days', '-d', type=int, default=14, help='過去何日以内のファイルを対象にするか指定（デフォルト: 14日）')
    parser.add_argument('--all', '-a', action='store_true', help='すべての期間のファイルを対象にする（--daysよりも優先）')
    parser.add_argument('--exclude', '-e', nargs='+', help='除外するフォルダやファイルのパターンを指定（スペース区切りで複数指定可能）')

    args = parser.parse_args()

    # 詳細ログが有効な場合はDEBUGレベルに設定
    if args.verbose:
        logger.setLevel(logging.DEBUG)
        # ファイルハンドラもDEBUGレベルに設定
        for handler in logger.handlers:
            if isinstance(handler, logging.FileHandler):
                handler.setLevel(logging.DEBUG)

    # フォルダパスの取得と検証
    target_folder_path_str = args.target_folder_path
    target_folder_path = Path(target_folder_path_str)

    if not target_folder_path.exists():
        logger.error(f"指定されたフォルダパスが見つかりません: {target_folder_path}")
        print(f"エラー: 指定されたフォルダパスが見つかりません: {target_folder_path}")
        return 1
    if not target_folder_path.is_dir():
        logger.error(f"指定されたパスはフォルダではありません: {target_folder_path}")
        print(f"エラー: 指定されたパスはフォルダではありません: {target_folder_path}")
        return 1

    # 期間指定の処理
    days = None if args.all else args.days

    # 除外パターンの処理
    exclude_patterns = args.exclude if args.exclude else []

    # output_dir の存在確認と出力ファイルパスの設定
    global output_dir
    if not output_dir.exists():
        logger.warning(f"出力ディレクトリが見つかりません: {output_dir}. 作成を試みます。")
        try:
            output_dir.mkdir(parents=True, exist_ok=True)
        except Exception as mkdir_e:
            logger.error(f"出力ディレクトリの作成に失敗しました: {output_dir}\nエラー: {mkdir_e}")
            print(f"エラー: 出力ディレクトリの作成に失敗しました: {output_dir}\nエラー: {mkdir_e}")
            return 1

    # 出力ファイル名に禁止文字が含まれていないか確認・置換
    folder_name = target_folder_path.name
    if not folder_name: # ルートディレクトリなどが指定された場合
        folder_name = "root_context"
    safe_folder_name = re.sub(r'[\\/*?:"<>|]', '_', folder_name)
    timestamp_str = datetime.datetime.now().strftime("%Y%m%d_%H%M%S") # タイムスタンプ追加
    output_filename = f"{safe_folder_name}_{timestamp_str}_contents.md" # タイムスタンプもファイル名に含める
    output_file = output_dir / output_filename

    logger.info(f"処理開始: フォルダパス '{target_folder_path}' の内容抽出")
    if days is not None:
        logger.info(f"期間制限: 過去 {days} 日以内")
    if exclude_patterns:
        logger.info(f"除外パターン: {exclude_patterns}")
    logger.info(f"出力ファイル: {output_file}")

    try:
        # フォルダとファイルの検索 (関数名と引数を変更)
        logger.info(f"ターゲットフォルダの絶対パス: {target_folder_path.resolve()}")
        found_folders_info = find_folders_and_files(target_folder_path, days, exclude_patterns)

        # Markdown ファイルの作成 (引数を変更)
        md_content = create_markdown_content(target_folder_path, found_folders_info)

        # ファイルに書き出し
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(md_content)

        logger.info(f"処理完了: 出力ファイル '{output_file}' を作成しました")
        print(f"\n出力ファイル: {output_file}")

    except Exception as e:
        logger.error(f"エラーが発生しました: {e}", exc_info=True)
        print(f"\nエラーが発生しました: {e}")
        return 1

    return 0

if __name__ == "__main__":
    sys.exit(main())
